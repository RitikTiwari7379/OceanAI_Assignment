from fastapi import FastAPI, APIRouter, Depends, HTTPException, status
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from fastapi.responses import FileResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
import os
import logging
from pathlib import Path
from pydantic import BaseModel, EmailStr
from typing import List, Optional, Dict, Any
import uuid
from datetime import datetime, timezone, timedelta
import aiosqlite
from contextlib import asynccontextmanager
import bcrypt
import jwt
import google.generativeai as genai
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
import tempfile
import json

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# Configuration
SECRET_KEY = os.environ.get('SECRET_KEY', 'your-secret-key-change-in-production')
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 8  # 8 hours (expires after work day)
DATABASE_PATH = ROOT_DIR / 'app.db'

# Gemini API setup
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
if GEMINI_API_KEY:
    print("✓ Gemini API Key loaded successfully")
    genai.configure(api_key=GEMINI_API_KEY)
else:
    print("❌ WARNING: GEMINI_API_KEY not found in environment variables")
    print("Please set GEMINI_API_KEY in your .env file")

@asynccontextmanager
async def lifespan(app: FastAPI):
    await init_db()
    yield

app = FastAPI(lifespan=lifespan)
api_router = APIRouter(prefix="/api")
security = HTTPBearer()

@app.get("/health")
async def health_check():
    """Health check endpoint for monitoring"""
    return {"status": "healthy", "service": "contentcraft-backend"}

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

async def init_db():
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Users table
        await db.execute("""
            CREATE TABLE IF NOT EXISTS users (
                id TEXT PRIMARY KEY,
                email TEXT UNIQUE NOT NULL,
                password_hash TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
        """)
        
        # Projects table
        await db.execute("""
            CREATE TABLE IF NOT EXISTS projects (
                id TEXT PRIMARY KEY,
                user_id TEXT NOT NULL,
                name TEXT NOT NULL,
                type TEXT NOT NULL,
                topic TEXT NOT NULL,
                config TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                FOREIGN KEY (user_id) REFERENCES users (id)
            )
        """)
        
        # Sections table
        await db.execute("""
            CREATE TABLE IF NOT EXISTS sections (
                id TEXT PRIMARY KEY,
                project_id TEXT NOT NULL,
                section_order INTEGER NOT NULL,
                title TEXT NOT NULL,
                content TEXT,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                FOREIGN KEY (project_id) REFERENCES projects (id)
            )
        """)
        
        # Refinements table
        await db.execute("""
            CREATE TABLE IF NOT EXISTS refinements (
                id TEXT PRIMARY KEY,
                section_id TEXT NOT NULL,
                prompt TEXT NOT NULL,
                response TEXT NOT NULL,
                feedback TEXT,
                comment TEXT,
                created_at TEXT NOT NULL,
                FOREIGN KEY (section_id) REFERENCES sections (id)
            )
        """)
        
        await db.commit()

class UserRegister(BaseModel):
    email: EmailStr
    password: str

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class TokenResponse(BaseModel):
    access_token: str
    token_type: str
    user_id: str
    email: str

class ProjectCreate(BaseModel):
    name: str
    type: str  # 'docx' or 'pptx'
    topic: str
    config: Dict[str, Any]

class ProjectResponse(BaseModel):
    id: str
    user_id: str
    name: str
    type: str
    topic: str
    config: Dict[str, Any]
    created_at: str
    updated_at: str

class SectionResponse(BaseModel):
    id: str
    project_id: str
    section_order: int
    title: str
    content: Optional[str]
    created_at: str
    updated_at: str

class GenerateContentRequest(BaseModel):
    project_id: str

class RefineContentRequest(BaseModel):
    section_id: str
    prompt: str

class FeedbackRequest(BaseModel):
    section_id: str
    feedback: str  # 'like' or 'dislike'
    comment: Optional[str] = None

class AITemplateRequest(BaseModel):
    type: str  # 'docx' or 'pptx'
    topic: str

# Auth utilities
def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_access_token(data: dict) -> str:
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)) -> dict:
    try:
        token = credentials.credentials
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        user_id = payload.get("sub")
        if user_id is None:
            raise HTTPException(status_code=401, detail="Invalid token")
        return {"user_id": user_id}
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except jwt.JWTError:
        raise HTTPException(status_code=401, detail="Invalid token")

@api_router.post("/auth/register", response_model=TokenResponse)
async def register(user: UserRegister):
    user_id = str(uuid.uuid4())
    password_hash = hash_password(user.password)
    created_at = datetime.now(timezone.utc).isoformat()
    
    async with aiosqlite.connect(DATABASE_PATH) as db:
        try:
            await db.execute(
                "INSERT INTO users (id, email, password_hash, created_at) VALUES (?, ?, ?, ?)",
                (user_id, user.email, password_hash, created_at)
            )
            await db.commit()
        except aiosqlite.IntegrityError:
            raise HTTPException(status_code=400, detail="Email already registered")
    
    access_token = create_access_token({"sub": user_id})
    return TokenResponse(
        access_token=access_token,
        token_type="bearer",
        user_id=user_id,
        email=user.email
    )

@api_router.post("/auth/login", response_model=TokenResponse)
async def login(user: UserLogin):
    async with aiosqlite.connect(DATABASE_PATH) as db:
        async with db.execute(
            "SELECT id, email, password_hash FROM users WHERE email = ?",
            (user.email,)
        ) as cursor:
            row = await cursor.fetchone()
            if not row or not verify_password(user.password, row[2]):
                raise HTTPException(status_code=401, detail="Invalid credentials")
            
            access_token = create_access_token({"sub": row[0]})
            return TokenResponse(
                access_token=access_token,
                token_type="bearer",
                user_id=row[0],
                email=row[1]
            )

@api_router.get("/auth/validate")
async def validate_token(current_user: dict = Depends(get_current_user)):
    """Validate if the current token is still valid and return user info"""
    async with aiosqlite.connect(DATABASE_PATH) as db:
        async with db.execute(
            "SELECT id, email FROM users WHERE id = ?",
            (current_user['user_id'],)
        ) as cursor:
            row = await cursor.fetchone()
            if not row:
                raise HTTPException(status_code=401, detail="User not found")
            
            return {
                "valid": True,
                "user_id": row[0],
                "email": row[1]
            }

# Project endpoints
@api_router.post("/projects", response_model=ProjectResponse)
async def create_project(project: ProjectCreate, current_user: dict = Depends(get_current_user)):
    project_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    async with aiosqlite.connect(DATABASE_PATH) as db:
        await db.execute(
            """INSERT INTO projects (id, user_id, name, type, topic, config, created_at, updated_at)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?)""",
            (project_id, current_user['user_id'], project.name, project.type, 
             project.topic, json.dumps(project.config), now, now)
        )
        await db.commit()
    
    return ProjectResponse(
        id=project_id,
        user_id=current_user['user_id'],
        name=project.name,
        type=project.type,
        topic=project.topic,
        config=project.config,
        created_at=now,
        updated_at=now
    )

@api_router.get("/projects", response_model=List[ProjectResponse])
async def get_projects(current_user: dict = Depends(get_current_user)):
    projects = []
    async with aiosqlite.connect(DATABASE_PATH) as db:
        async with db.execute(
            "SELECT id, user_id, name, type, topic, config, created_at, updated_at FROM projects WHERE user_id = ? ORDER BY created_at DESC",
            (current_user['user_id'],)
        ) as cursor:
            async for row in cursor:
                projects.append(ProjectResponse(
                    id=row[0],
                    user_id=row[1],
                    name=row[2],
                    type=row[3],
                    topic=row[4],
                    config=json.loads(row[5]),
                    created_at=row[6],
                    updated_at=row[7]
                ))
    return projects

@api_router.get("/projects/{project_id}", response_model=ProjectResponse)
async def get_project(project_id: str, current_user: dict = Depends(get_current_user)):
    async with aiosqlite.connect(DATABASE_PATH) as db:
        async with db.execute(
            "SELECT id, user_id, name, type, topic, config, created_at, updated_at FROM projects WHERE id = ? AND user_id = ?",
            (project_id, current_user['user_id'])
        ) as cursor:
            row = await cursor.fetchone()
            if not row:
                raise HTTPException(status_code=404, detail="Project not found")
            return ProjectResponse(
                id=row[0],
                user_id=row[1],
                name=row[2],
                type=row[3],
                topic=row[4],
                config=json.loads(row[5]),
                created_at=row[6],
                updated_at=row[7]
            )

@api_router.delete("/projects/{project_id}")
async def delete_project(project_id: str, current_user: dict = Depends(get_current_user)):
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Verify project ownership
        async with db.execute(
            "SELECT id FROM projects WHERE id = ? AND user_id = ?",
            (project_id, current_user['user_id'])
        ) as cursor:
            if not await cursor.fetchone():
                raise HTTPException(status_code=404, detail="Project not found")
        
        # Delete refinements first (foreign key constraint)
        await db.execute("""
            DELETE FROM refinements 
            WHERE section_id IN (
                SELECT id FROM sections WHERE project_id = ?
            )
        """, (project_id,))
        
        # Delete sections
        await db.execute("DELETE FROM sections WHERE project_id = ?", (project_id,))
        
        # Delete project
        await db.execute("DELETE FROM projects WHERE id = ?", (project_id,))
        
        await db.commit()
        return {"message": "Project deleted successfully"}

async def generate_with_gemini(prompt: str) -> str:
    try:
        model = genai.GenerativeModel('gemini-2.5-flash')
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        logger.error(f"Gemini API error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"AI generation failed: {str(e)}")

@api_router.post("/generate-content")
async def generate_content(request: GenerateContentRequest, current_user: dict = Depends(get_current_user)):
    # Get project
    async with aiosqlite.connect(DATABASE_PATH) as db:
        async with db.execute(
            "SELECT id, type, topic, config FROM projects WHERE id = ? AND user_id = ?",
            (request.project_id, current_user['user_id'])
        ) as cursor:
            project_row = await cursor.fetchone()
            if not project_row:
                raise HTTPException(status_code=404, detail="Project not found")
        
        project_id, doc_type, topic, config_str = project_row
        config = json.loads(config_str)
        
        # Check if content already generated or is being generated
        async with db.execute(
            "SELECT COUNT(*) FROM sections WHERE project_id = ?",
            (project_id,)
        ) as cursor:
            count = await cursor.fetchone()
            if count[0] > 0:
                # Content already exists, return existing content
                sections = []
                async with db.execute(
                    "SELECT id, title, content, section_order FROM sections WHERE project_id = ? ORDER BY section_order",
                    (project_id,)
                ) as cursor:
                    async for row in cursor:
                        sections.append({
                            'id': row[0],
                            'title': row[1], 
                            'content': row[2],
                            'order': row[3]
                        })
                return {"message": "Content already generated", "sections": sections}
        
        # Generate content for each section
        sections = []
        now = datetime.now(timezone.utc).isoformat()
        
        if doc_type == 'docx':
            outline = config.get('outline', [])
            for i, section_title in enumerate(outline):
                section_id = str(uuid.uuid4())
                prompt = f"""Write detailed, professional content for a document section.

Document Topic: {topic}
Section Title: {section_title}
Section Number: {i+1} of {len(outline)}

Provide comprehensive content (300-500 words) for this section. Focus on being informative, well-structured, and professional. Do not include the section title in your response."""
                
                content = await generate_with_gemini(prompt)
                
                await db.execute(
                    """INSERT INTO sections (id, project_id, section_order, title, content, created_at, updated_at)
                       VALUES (?, ?, ?, ?, ?, ?, ?)""",
                    (section_id, project_id, i, section_title, content, now, now)
                )
                sections.append({
                    'id': section_id,
                    'title': section_title,
                    'content': content,
                    'order': i
                })
        
        elif doc_type == 'pptx':
            slides = config.get('slides', [])
            for i, slide_title in enumerate(slides):
                section_id = str(uuid.uuid4())
                prompt = f"""Create content for a PowerPoint slide.

Presentation Topic: {topic}
Slide Title: {slide_title}
Slide Number: {i+1} of {len(slides)}

Provide 3-5 concise bullet points for this slide. Each point should be clear and impactful. Format as a simple bulleted list. Do not include the slide title."""
                
                content = await generate_with_gemini(prompt)
                
                await db.execute(
                    """INSERT INTO sections (id, project_id, section_order, title, content, created_at, updated_at)
                       VALUES (?, ?, ?, ?, ?, ?, ?)""",
                    (section_id, project_id, i, slide_title, content, now, now)
                )
                sections.append({
                    'id': section_id,
                    'title': slide_title,
                    'content': content,
                    'order': i
                })
        
        await db.commit()
        logger.info(f"Successfully generated {len(sections)} sections for project {project_id}")
        return {"message": "Content generated successfully", "sections": sections}

@api_router.get("/projects/{project_id}/sections", response_model=List[SectionResponse])
async def get_sections(project_id: str, current_user: dict = Depends(get_current_user)):
    sections = []
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Verify project ownership
        async with db.execute(
            "SELECT id FROM projects WHERE id = ? AND user_id = ?",
            (project_id, current_user['user_id'])
        ) as cursor:
            if not await cursor.fetchone():
                raise HTTPException(status_code=404, detail="Project not found")
        
        # Get sections
        async with db.execute(
            "SELECT id, project_id, section_order, title, content, created_at, updated_at FROM sections WHERE project_id = ? ORDER BY section_order",
            (project_id,)
        ) as cursor:
            async for row in cursor:
                sections.append(SectionResponse(
                    id=row[0],
                    project_id=row[1],
                    section_order=row[2],
                    title=row[3],
                    content=row[4],
                    created_at=row[5],
                    updated_at=row[6]
                ))
    return sections

@api_router.post("/refine-content")
async def refine_content(request: RefineContentRequest, current_user: dict = Depends(get_current_user)):
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Get section and verify ownership
        async with db.execute(
            """SELECT s.id, s.content, s.title, p.user_id, p.topic 
               FROM sections s 
               JOIN projects p ON s.project_id = p.id 
               WHERE s.id = ?""",
            (request.section_id,)
        ) as cursor:
            row = await cursor.fetchone()
            if not row:
                raise HTTPException(status_code=404, detail="Section not found")
            if row[3] != current_user['user_id']:
                raise HTTPException(status_code=403, detail="Not authorized")
            
            section_id, current_content, section_title, _, topic = row
        
        # Generate refined content
        prompt = f"""You are refining content for a section.

Section Title: {section_title}
Document Topic: {topic}
Current Content:
{current_content}

User Request: {request.prompt}

Provide the refined version of the content based on the user's request. Return only the refined content, no explanations."""
        
        refined_content = await generate_with_gemini(prompt)
        
        # Update section
        now = datetime.now(timezone.utc).isoformat()
        await db.execute(
            "UPDATE sections SET content = ?, updated_at = ? WHERE id = ?",
            (refined_content, now, section_id)
        )
        
        # Save refinement history
        refinement_id = str(uuid.uuid4())
        await db.execute(
            """INSERT INTO refinements (id, section_id, prompt, response, created_at)
               VALUES (?, ?, ?, ?, ?)""",
            (refinement_id, section_id, request.prompt, refined_content, now)
        )
        
        await db.commit()
        
        return {
            "message": "Content refined successfully",
            "section_id": section_id,
            "content": refined_content
        }

@api_router.post("/feedback")
async def add_feedback(request: FeedbackRequest, current_user: dict = Depends(get_current_user)):
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Verify section ownership
        async with db.execute(
            """SELECT s.id FROM sections s 
               JOIN projects p ON s.project_id = p.id 
               WHERE s.id = ? AND p.user_id = ?""",
            (request.section_id, current_user['user_id'])
        ) as cursor:
            if not await cursor.fetchone():
                raise HTTPException(status_code=404, detail="Section not found")
        
        # Get latest refinement or create new one
        async with db.execute(
            "SELECT id FROM refinements WHERE section_id = ? ORDER BY created_at DESC LIMIT 1",
            (request.section_id,)
        ) as cursor:
            row = await cursor.fetchone()
            
            if row:
                # Update existing refinement
                await db.execute(
                    "UPDATE refinements SET feedback = ?, comment = ? WHERE id = ?",
                    (request.feedback, request.comment, row[0])
                )
            else:
                # Create new refinement record for feedback
                refinement_id = str(uuid.uuid4())
                now = datetime.now(timezone.utc).isoformat()
                await db.execute(
                    """INSERT INTO refinements (id, section_id, prompt, response, feedback, comment, created_at)
                       VALUES (?, ?, ?, ?, ?, ?, ?)""",
                    (refinement_id, request.section_id, 'Initial feedback', '', request.feedback, request.comment, now)
                )
            
            await db.commit()
            return {"message": "Feedback saved successfully"}

@api_router.get("/sections/{section_id}/feedback")
async def get_section_feedback(section_id: str, current_user: dict = Depends(get_current_user)):
    """Get feedback (like/dislike) state for a specific section"""
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Verify section belongs to user's project
        async with db.execute(
            """SELECT s.id FROM sections s
               JOIN projects p ON s.project_id = p.id 
               WHERE s.id = ? AND p.user_id = ?""",
            (section_id, current_user['user_id'])
        ) as cursor:
            if not await cursor.fetchone():
                raise HTTPException(status_code=404, detail="Section not found")
        
        # Get latest feedback from refinements table
        async with db.execute(
            "SELECT feedback, comment FROM refinements WHERE section_id = ? AND feedback IS NOT NULL ORDER BY created_at DESC LIMIT 1",
            (section_id,)
        ) as cursor:
            row = await cursor.fetchone()
            
        if row and row[0]:
            return {
                "feedback": row[0],
                "comment": row[1] or "",
                "liked": row[0] == "like",
                "disliked": row[0] == "dislike"
            }
        else:
            return {
                "feedback": None,
                "comment": "",
                "liked": False,
                "disliked": False
            }

@api_router.get("/sections/{section_id}/comments")
async def get_section_comments(section_id: str, current_user: dict = Depends(get_current_user)):
    """Get all comments for a specific section"""
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Verify section belongs to user's project
        async with db.execute(
            """SELECT s.id FROM sections s
               JOIN projects p ON s.project_id = p.id 
               WHERE s.id = ? AND p.user_id = ?""",
            (section_id, current_user['user_id'])
        ) as cursor:
            if not await cursor.fetchone():
                raise HTTPException(status_code=404, detail="Section not found")
        
        # Get all comments from refinements table
        async with db.execute(
            "SELECT comment, created_at FROM refinements WHERE section_id = ? AND comment IS NOT NULL AND comment != '' ORDER BY created_at DESC",
            (section_id,)
        ) as cursor:
            comments = await cursor.fetchall()
            
        return [{"comment": row[0], "created_at": row[1]} for row in comments]

@api_router.get("/sections/{section_id}/revisions")
async def get_section_revisions(section_id: str, current_user: dict = Depends(get_current_user)):
    """Get all revision history (prompts and responses) for a specific section"""
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Verify section belongs to user's project
        async with db.execute(
            """SELECT s.id FROM sections s
               JOIN projects p ON s.project_id = p.id 
               WHERE s.id = ? AND p.user_id = ?""",
            (section_id, current_user['user_id'])
        ) as cursor:
            if not await cursor.fetchone():
                raise HTTPException(status_code=404, detail="Section not found")
        
        # Get all refinements (excluding initial feedback records)
        async with db.execute(
            "SELECT prompt, response, created_at FROM refinements WHERE section_id = ? AND prompt != 'Initial feedback' ORDER BY created_at DESC",
            (section_id,)
        ) as cursor:
            revisions = await cursor.fetchall()
            
        return [
            {
                "prompt": row[0], 
                "response": row[1], 
                "created_at": row[2],
                "timestamp": row[2]
            } for row in revisions
        ]

@api_router.post("/ai-template", response_model=Dict[str, Any])
async def generate_ai_template(request: AITemplateRequest, current_user: dict = Depends(get_current_user)):
    if request.type == 'docx':
        prompt = f"""Generate a document outline for the following topic: {request.topic}

Provide 5-8 section headings that would make a comprehensive, well-structured document. Return only the section headings, one per line, no numbering or additional text."""
    else:  # pptx
        prompt = f"""Generate slide titles for a PowerPoint presentation on: {request.topic}

Provide 8-12 slide titles that would make a comprehensive presentation. Include an introduction slide and a conclusion slide. Return only the slide titles, one per line, no numbering."""
    
    try:
        response = await generate_with_gemini(prompt)
        items = [line.strip() for line in response.strip().split('\n') if line.strip()]
        
        return {
            'type': request.type,
            'topic': request.topic,
            'items': items
        }
    except Exception as e:
        logger.error(f"AI template generation error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/export/{project_id}")
async def export_document(project_id: str, current_user: dict = Depends(get_current_user)):
    async with aiosqlite.connect(DATABASE_PATH) as db:
        # Get project
        async with db.execute(
            "SELECT name, type, topic FROM projects WHERE id = ? AND user_id = ?",
            (project_id, current_user['user_id'])
        ) as cursor:
            project_row = await cursor.fetchone()
            if not project_row:
                raise HTTPException(status_code=404, detail="Project not found")
            
            project_name, doc_type, topic = project_row
        
        # Get sections
        sections = []
        async with db.execute(
            "SELECT title, content FROM sections WHERE project_id = ? ORDER BY section_order",
            (project_id,)
        ) as cursor:
            async for row in cursor:
                sections.append({'title': row[0], 'content': row[1]})
        
        if not sections:
            raise HTTPException(status_code=400, detail="No content to export")
        
        # Generate document
        if doc_type == 'docx':
            doc = Document()
            
            # Add title
            title = doc.add_heading(project_name, 0)
            title.alignment = 1  # Center
            
            # Add sections
            for section in sections:
                doc.add_heading(section['title'], 1)
                doc.add_paragraph(section['content'] or '')
            
            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
            doc.save(temp_file.name)
            temp_file.close()
            
            return FileResponse(
                temp_file.name,
                media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                filename=f"{project_name.replace(' ', '_')}.docx"
            )
        
        elif doc_type == 'pptx':
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = project_name
            subtitle.text = topic
            
            # Content slides
            for section in sections:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = section['title']
                
                tf = body_shape.text_frame
                content = section['content'] or ''
                
                # Parse content into bullet points
                lines = [line.strip() for line in content.split('\n') if line.strip()]
                for line in lines:
                    # Remove bullet characters if present
                    clean_line = line.lstrip('•-*').strip()
                    if clean_line:
                        p = tf.add_paragraph()
                        p.text = clean_line
                        p.level = 0
            
            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
            prs.save(temp_file.name)
            temp_file.close()
            
            return FileResponse(
                temp_file.name,
                media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                filename=f"{project_name.replace(' ', '_')}.pptx"
            )

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
